import importlib
import io
import os
from functools import wraps
from typing import Any


def with_input_output_key(fn):
    """Decorator that wraps a parser function that takes a single
    string parameter and return list of strings and makes it a full
    parser function that takes an item as a dictionary and return a
    list of dictionaries."""

    @wraps(fn)
    def wrapper(item, input_key="text", output_key="text", **kw):
        if input_key not in item:
            raise ValueError(f"Input key {input_key} not found in item: {item}")
        result = fn(item[input_key], **kw)
        if not isinstance(result, list):
            result = [result]
        return [{output_key: res} for res in result]

    return wrapper


def llama_index_simple_directory_reader(
    item: dict[str, Any], input_key: str = "path"
) -> list[dict[str, Any]]:
    from llama_index.core import SimpleDirectoryReader

    documents = SimpleDirectoryReader(item[input_key]).load_data()
    return [{"text": doc.text, "metadata": doc.metadata} for doc in documents]


def llama_index_wikipedia_reader(
    item: dict[str, Any], input_key: str = "pages"
) -> list[dict[str, Any]]:
    from llama_index.readers.wikipedia import WikipediaReader

    loader = WikipediaReader()
    pages = item[input_key]
    if not isinstance(pages, list):
        pages = [pages]
    documents = loader.load_data(pages=pages, auto_suggest=False)
    # The wikipedia reader does not include the page url in the metadata, which is impractical...
    for name, doc in zip(pages, documents):
        doc.metadata["source"] = "https://en.wikipedia.org/wiki/" + name

    return [{"text": doc.text, "metadata": doc.metadata} for doc in documents]


@with_input_output_key
def whisper_speech_to_text(filename: str) -> list[str]:
    """
    Transcribe speech from an audio file to text using Whisper model via litellm.
    If the file is larger than 25 MB, it's split into 10-minute chunks with 30-second overlap.

    Args:
        filename (str): Path to the mp3 or mp4 file.

    Returns:
        list[str]: Transcribed text.
    """

    from litellm import transcription

    file_size = os.path.getsize(filename)
    if file_size > 25 * 1024 * 1024:  # 25 MB in bytes
        from pydub import AudioSegment

        audio = AudioSegment.from_file(filename)
        chunk_length = 10 * 60 * 1000  # 10 minutes in milliseconds
        overlap = 30 * 1000  # 30 seconds in milliseconds

        chunks = []
        for i in range(0, len(audio), chunk_length - overlap):
            chunk = audio[i : i + chunk_length]
            chunks.append(chunk)

        transcriptions = []

        for i, chunk in enumerate(chunks):
            buffer = io.BytesIO()
            buffer.name = f"temp_chunk_{i}_{os.path.basename(filename)}"
            chunk.export(buffer, format="mp3")
            buffer.seek(0)  # Reset buffer position to the beginning

            response = transcription(model="whisper-1", file=buffer)
            transcriptions.append(response.text)

        return transcriptions
    else:
        with open(filename, "rb") as audio_file:
            response = transcription(model="whisper-1", file=audio_file)

        return [response.text]


@with_input_output_key
def xlsx_to_string(
    filename: str,
    orientation: str = "col",
    col_order: list[str] | None = None,
    doc_per_sheet: bool = False,
) -> list[str]:
    """
    Convert an Excel file to a string representation or a list of string representations.

    Args:
        filename (str): Path to the xlsx file.
        orientation (str): Either "row" or "col" for cell arrangement.
        col_order (list[str] | None): List of column names to specify the order.
        doc_per_sheet (bool): If True, return a list of strings, one per sheet.

    Returns:
        list[str]: String representation(s) of the Excel file content.
    """
    import openpyxl

    wb = openpyxl.load_workbook(filename)

    def process_sheet(sheet):
        if col_order:
            headers = [
                col for col in col_order if col in sheet.iter_cols(1, sheet.max_column)
            ]
        else:
            headers = [cell.value for cell in sheet[1]]

        result = []
        if orientation == "col":
            for col_idx, header in enumerate(headers, start=1):
                column = sheet.cell(1, col_idx).column_letter
                column_values = [cell.value for cell in sheet[column][1:]]
                result.append(f"{header}: " + "\n".join(map(str, column_values)))
                result.append("")  # Empty line between columns
        else:  # row
            for row in sheet.iter_rows(min_row=2, values_only=True):
                row_dict = {
                    header: value for header, value in zip(headers, row) if header
                }
                result.append(
                    " | ".join(
                        [f"{header}: {value}" for header, value in row_dict.items()]
                    )
                )

        return "\n".join(result)

    if doc_per_sheet:
        return [process_sheet(sheet) for sheet in wb.worksheets]
    else:
        return [process_sheet(wb.active)]


@with_input_output_key
def txt_to_string(filename: str) -> list[str]:
    """
    Read the content of a text file and return it as a list of strings (only one element).

    Args:
        filename (str): Path to the txt or md file.

    Returns:
        list[str]: Content of the file as a list of strings.
    """
    with open(filename, "r", encoding="utf-8") as file:
        return [file.read()]


@with_input_output_key
def docx_to_string(filename: str) -> list[str]:
    """
    Extract text from a Word document.

    Args:
        filename (str): Path to the docx file.

    Returns:
        list[str]: Extracted text from the document.
    """
    from docx import Document

    doc = Document(filename)
    return ["\n".join([paragraph.text for paragraph in doc.paragraphs])]


@with_input_output_key
def pptx_to_string(filename: str, doc_per_slide: bool = False) -> list[str]:
    """
    Extract text from a PowerPoint presentation.

    Args:
        filename (str): Path to the pptx file.
        doc_per_slide (bool): If True, return each slide as a separate
            document. If False, return the entire presentation as one document.

    Returns:
        list[str]: Extracted text from the presentation. If doc_per_slide
            is True, each string in the list represents a single slide.
            Otherwise, the list contains a single string with all slides'
            content.
    """
    from pptx import Presentation

    prs = Presentation(filename)
    result = []

    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)

        if doc_per_slide:
            result.append("\n".join(slide_content))
        else:
            result.extend(slide_content)

    if not doc_per_slide:
        result = ["\n".join(result)]

    return result


@with_input_output_key
def azure_di_read(
    filename: str,
    use_url: bool = False,
    include_line_numbers: bool = False,
    include_handwritten: bool = False,
    include_font_styles: bool = False,
    include_selection_marks: bool = False,
    doc_per_page: bool = False,
) -> list[str]:
    """
    > Note to developers: We used [this documentation](https://learn.microsoft.com/en-us/azure/ai-services/document-intelligence/how-to-guides/use-sdk-rest-api?view=doc-intel-4.0.0&tabs=windows&pivots=programming-language-python) as a reference.

    This function uses Azure Document Intelligence to extract text from documents.
    To use this function, you need to set up an Azure Document Intelligence resource:

    1. [Create an Azure account](https://azure.microsoft.com/) if you don't have one
    2. Set up a Document Intelligence resource in the [Azure portal](https://portal.azure.com/#create/Microsoft.CognitiveServicesFormRecognizer)
    3. Once created, find the resource's endpoint and key in the Azure portal
    4. Set these as environment variables:
       - DOCUMENTINTELLIGENCE_API_KEY: Your Azure Document Intelligence API key
       - DOCUMENTINTELLIGENCE_ENDPOINT: Your Azure Document Intelligence endpoint URL

    The function will use these credentials to authenticate with the Azure service.
    If the environment variables are not set, the function will raise a ValueError.

    The Azure Document Intelligence client is then initialized with these credentials.
    It sends the document (either as a file or URL) to Azure for analysis.
    The service processes the document and returns structured information about its content.

    This function then extracts the text content from the returned data,
    applying any specified formatting options (like including line numbers or font styles).
    The extracted text is returned as a list of strings, with each string
    representing either a page (if doc_per_page is True) or the entire document.

    Args:
        filename (str): Path to the file to be analyzed or URL of the document if use_url is True.
        use_url (bool, optional): If True, treat filename as a URL. Defaults to False.
        include_line_numbers (bool, optional): If True, include line numbers in the output. Defaults to False.
        include_handwritten (bool, optional): If True, include handwritten text in the output. Defaults to False.
        include_font_styles (bool, optional): If True, include font style information in the output. Defaults to False.
        include_selection_marks (bool, optional): If True, include selection marks in the output. Defaults to False.
        doc_per_page (bool, optional): If True, return each page as a separate document. Defaults to False.

    Returns:
        list[str]: Extracted text from the document. If doc_per_page is True, each string in the list represents
                   a single page. Otherwise, the list contains a single string with all pages' content.

    Raises:
        ValueError: If DOCUMENTINTELLIGENCE_API_KEY or DOCUMENTINTELLIGENCE_ENDPOINT environment variables are not set.
    """

    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
    from azure.core.credentials import AzureKeyCredential

    key = os.getenv("DOCUMENTINTELLIGENCE_API_KEY")
    endpoint = os.getenv("DOCUMENTINTELLIGENCE_ENDPOINT")

    if key is None:
        raise ValueError("DOCUMENTINTELLIGENCE_API_KEY environment variable is not set")
    if endpoint is None:
        raise ValueError(
            "DOCUMENTINTELLIGENCE_ENDPOINT environment variable is not set"
        )

    document_analysis_client = DocumentIntelligenceClient(
        endpoint=endpoint, credential=AzureKeyCredential(key)
    )

    if use_url:
        poller = document_analysis_client.begin_analyze_document(
            "prebuilt-read", AnalyzeDocumentRequest(url_source=filename)
        )
    else:
        with open(filename, "rb") as f:
            poller = document_analysis_client.begin_analyze_document("prebuilt-read", f)

    result = poller.result()

    style_content = []
    content = []

    if result.styles:
        for style in result.styles:
            if style.is_handwritten and include_handwritten:
                handwritten_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"Handwritten content: {handwritten_text}")

            if style.font_style and include_font_styles:
                styled_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"'{style.font_style}' font style: {styled_text}")

    for page in result.pages:
        page_content = []

        if page.lines:
            for line_idx, line in enumerate(page.lines):
                if include_line_numbers:
                    page_content.append(f" Line #{line_idx}: {line.content}")
                else:
                    page_content.append(f"{line.content}")

        if page.selection_marks and include_selection_marks:
            # TODO: figure this out
            for selection_mark_idx, selection_mark in enumerate(page.selection_marks):
                page_content.append(
                    f"Selection mark #{selection_mark_idx}: State is '{selection_mark.state}' within bounding polygon "
                    f"'{selection_mark.polygon}' and has a confidence of {selection_mark.confidence}"
                )

        content.append("\n".join(page_content))

    if doc_per_page:
        return style_content + content
    else:
        return [
            "\n\n".join(
                [
                    "\n".join(style_content),
                    "\n\n".join(
                        f"Page {i+1}:\n{page_content}"
                        for i, page_content in enumerate(content)
                    ),
                ]
            )
        ]


@with_input_output_key
def paddleocr_pdf_to_string(
    input_path: str,
    doc_per_page: bool = False,
    ocr_enabled: bool = True,
    lang: str = "en",
) -> list[str]:
    """
    Extract text and image information from a PDF file using PaddleOCR for image-based PDFs.

    **Note: this is very slow!!**

    Args:
        input_path (str): Path to the input PDF file.
        doc_per_page (bool): If True, return a list of strings, one per page.
            If False, return a single string.
        ocr_enabled (bool): Whether to enable OCR for image-based PDFs.
        lang (str): Language of the PDF file.

    Returns:
        list[str]: Extracted content as a list of formatted strings.
    """
    import fitz
    import numpy as np
    from paddleocr import PaddleOCR

    ocr = PaddleOCR(use_angle_cls=True, lang=lang)

    pdf_content = []

    with fitz.open(input_path) as pdf:
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            text = page.get_text()
            images = []

            # Extract image information
            for img_index, img in enumerate(page.get_images(full=True)):
                rect = page.get_image_bbox(img)
                images.append(f"Image {img_index + 1}: bbox {rect}")

            page_content = f"Page {page_num + 1}:\n"
            page_content += f"Text:\n{text}\n"
            page_content += "Images:\n" + "\n".join(images) + "\n"

            if not text and ocr_enabled:
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                    pix.height, pix.width, 3
                )

                ocr_result = ocr.ocr(img, cls=True)
                page_content += "OCR Results:\n"
                for line in ocr_result[0]:
                    bbox, (text, _) = line
                    page_content += f"{bbox}, {text}\n"

            pdf_content.append(page_content)

    if not doc_per_page:
        return ["\n\n".join(pdf_content)]

    return pdf_content


@with_input_output_key
def gptpdf_to_string(
    input_path: str,
    gpt_model: str,
    api_key: str,
    base_url: str,
    verbose: bool = False,
    custom_prompt: dict[str, str] | None = None,
) -> str:
    """
    Parse PDF using GPT to convert the content of a PDF to a markdown format and write it to an output file.

    **Note: pip install gptpdf required**

    Args:
        input_path (str): Path to the input PDF file.
        gpt_model (str): GPT model to be used for parsing.
        api_key (str): API key for GPT service.
        base_url (str): Base URL for the GPT service.
        verbose (bool): If True, will print additional information during parsing.
        custom_prompt (dict[str, str] | None): Custom prompt for the GPT model. See https://github.com/CosmosShadow/gptpdf for more information.

    Returns:
        list[str]: Extracted content as a list of strings.
    """
    import tempfile

    from gptpdf import parse_pdf

    with tempfile.TemporaryDirectory() as temp_dir:
        kwargs = {
            "pdf_path": input_path,
            "output_dir": temp_dir,
            "api_key": api_key,
            "base_url": base_url,
            "model": gpt_model,
            "verbose": verbose,
        }
        if custom_prompt:
            kwargs["prompt"] = custom_prompt

        parsed_content, _ = parse_pdf(
            **kwargs
        )  # The second element is a list of image paths, which we don't need.

        return [parsed_content]


# Define a dictionary mapping function names to their corresponding functions


def get_parser(name: str):
    try:
        entrypoint = importlib.metadata.entry_points(group="docetl.parser")[name]
    except KeyError:
        raise KeyError(f"Unrecognized parser {name}")
    return entrypoint.load()


def get_parsing_tools():
    return [ep.name for ep in importlib.metadata.entry_points(group="docetl.parser")]
