import pytest
import os
import tempfile
from docetl import parsing_tools


@pytest.fixture
def temp_audio_file():
    import requests

    url = "https://listenaminute.com/a/animals.mp3"
    response = requests.get(url)
    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_file:
        temp_file.write(response.content)
    yield temp_file.name
    return temp_file.name


@pytest.fixture
def temp_xlsx_file():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "New York"])
    ws.append(["Bob", 25, "London"])
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as temp_file:
        wb.save(temp_file.name)
    yield temp_file.name
    return temp_file.name


@pytest.fixture
def temp_txt_file():
    with tempfile.NamedTemporaryFile(
        mode="w", suffix=".txt", delete=False
    ) as temp_file:
        temp_file.write("This is a test text file.\nIt has multiple lines.")
    yield temp_file.name
    return temp_file.name


@pytest.fixture
def temp_docx_file():
    from docx import Document

    doc = Document()
    doc.add_paragraph("This is a test Word document.")
    doc.add_paragraph("It has multiple paragraphs.")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
        doc.save(temp_file.name)
    yield temp_file.name
    return temp_file.name


@pytest.fixture
def temp_pptx_file():
    from pptx import Presentation

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Test Presentation"
    slide1.placeholders[1].text = "This is the first slide"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "This is the second slide"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        prs.save(temp_file.name)
    yield temp_file.name
    return temp_file.name


def test_whisper_speech_to_text(temp_audio_file):
    result = parsing_tools.whisper_speech_to_text.__wrapped__(temp_audio_file)

    assert isinstance(result, list)
    assert len(result) == 1
    assert isinstance(result[0], str)
    assert len(result[0]) > 0  # Ensure some text was transcribed


def test_xlsx_to_string(temp_xlsx_file):
    result = parsing_tools.xlsx_to_string.__wrapped__(temp_xlsx_file)

    assert isinstance(result, list)
    assert len(result) == 1
    assert "Name: Alice" in result[0]
    assert "Age: 30" in result[0]
    assert "City: New York" in result[0]


def test_xlsx_to_string_row_orientation(temp_xlsx_file):
    result = parsing_tools.xlsx_to_string.__wrapped__(temp_xlsx_file, orientation="row")

    assert isinstance(result, list)
    assert len(result) == 1
    assert "Name: Alice | Age: 30 | City: New York" in result[0]
    assert "Name: Bob | Age: 25 | City: London" in result[0]


def test_xlsx_to_string_doc_per_sheet(temp_xlsx_file):
    result = parsing_tools.xlsx_to_string.__wrapped__(
        temp_xlsx_file, doc_per_sheet=True
    )

    assert isinstance(result, list)
    assert len(result) == 1  # Only one sheet in our test file
    assert "Name: Alice" in result[0]
    assert "Age: 30" in result[0]
    assert "City: New York" in result[0]


def test_txt_to_string(temp_txt_file):
    result = parsing_tools.txt_to_string.__wrapped__(temp_txt_file)

    assert isinstance(result, list)
    assert len(result) == 1
    assert result[0] == "This is a test text file.\nIt has multiple lines."


def test_docx_to_string(temp_docx_file):
    result = parsing_tools.docx_to_string.__wrapped__(temp_docx_file)

    assert isinstance(result, list)
    assert len(result) == 1
    assert "This is a test Word document." in result[0]
    assert "It has multiple paragraphs." in result[0]


def test_pptx_to_string(temp_pptx_file):
    result = parsing_tools.pptx_to_string.__wrapped__(temp_pptx_file)

    assert isinstance(result, list)
    assert len(result) == 1
    assert "Test Presentation" in result[0]
    assert "This is the first slide" in result[0]
    assert "Second Slide" in result[0]
    assert "This is the second slide" in result[0]


def test_pptx_to_string_doc_per_slide(temp_pptx_file):
    result = parsing_tools.pptx_to_string.__wrapped__(
        temp_pptx_file, doc_per_slide=True
    )

    assert isinstance(result, list)
    assert len(result) == 2
    assert "Test Presentation" in result[0]
    assert "This is the first slide" in result[0]
    assert "Second Slide" in result[1]
    assert "This is the second slide" in result[1]


@pytest.fixture
def pdf_url():
    return "https://raw.githubusercontent.com/Azure-Samples/cognitive-services-REST-api-samples/master/curl/form-recognizer/rest-api/read.png"


def test_azure_di_read(pdf_url):
    # Test with default parameters
    result = parsing_tools.azure_di_read.__wrapped__(pdf_url, use_url=True)

    assert isinstance(result, list)
    assert len(result) == 1
    content = result[0]

    # Check for expected content in the extracted text
    assert "While healthcare is still in the early stages of its" in content
    assert "seeing pharmaceutical and other life sciences organizations" in content
    assert "Enhancing the patient" in content

    # Test with include_line_numbers=True
    result_line_numbers = parsing_tools.azure_di_read.__wrapped__(
        pdf_url, use_url=True, include_line_numbers=True
    )
    assert isinstance(result_line_numbers, list)
    assert len(result_line_numbers) == 1
    assert any("Line #" in line for line in result_line_numbers[0].split("\n"))


def test_azure_di_read_invoice():
    invoice_url = "https://raw.githubusercontent.com/Azure-Samples/cognitive-services-REST-api-samples/master/curl/form-recognizer/sample-invoice.pdf"
    result = parsing_tools.azure_di_read.__wrapped__(invoice_url, use_url=True)

    assert isinstance(result, list)
    assert len(result) == 1
    content = result[0]

    # Check for expected content in the extracted text
    assert "Contoso" in content


# Clean up temporary files after all tests have passed
def pytest_sessionfinish(session, exitstatus):
    if exitstatus == 0:
        for fixture in [
            temp_audio_file,
            temp_xlsx_file,
            temp_txt_file,
            temp_docx_file,
            temp_pptx_file,
        ]:
            file_path = session.config.cache.get(fixture.__name__, None)
            if file_path and os.path.exists(file_path):
                os.remove(file_path)
                os.unlink(file_path)


def test_paddleocr_pdf_to_string():
    pdf_path = "tests/data/PublicWaterMassMailing.pdf"
    result = parsing_tools.paddleocr_pdf_to_string.__wrapped__(pdf_path, lang="en")

    assert isinstance(result, list)
    assert len(result) == 1

    assert "have received the new bottles, please discard" in result[0]


# test function todo


def test_gptpdf_to_string():
    input_pdf = "tests/data/PublicWaterMassMailing.pdf"

    result = parsing_tools.gptpdf_to_string.__wrapped__(
        input_path=input_pdf,
        gpt_model="gpt-4o-mini",
        api_key=os.environ["OPENAI_API_KEY"],
        base_url="https://api.openai.com/v1",
        verbose=False,
    )

    assert len(result) > 0, "The extracted content should not be empty."
    assert len(result[0]) > 0, "The extracted content should not be empty."
